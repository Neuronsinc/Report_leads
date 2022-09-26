from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from datetime import date
from pptx.enum.shapes import MSO_SHAPE

SLD_LAYOUT_TITLE_AND_CONTENT = 1

EFECTIVIDAD_BAJA = RGBColor(0xFF, 0x00, 0x00)
EFECTIVIDAD_MEDIA = RGBColor(0xFF,0xA5,0x00)
EFECTIVIDAD_ALTA = RGBColor(0x00, 0xFF, 0x00)

leads_citas = [
    [100,200,300],
    [50,60,10],
    [5,2,3]
]


dates = ['01/09/2022 - 08/09/2022', '08/09/2022 - 15/09/2022', '15/09/2022 - 22/09/2022']

efectividades = [
    [
        ['30.48%', EFECTIVIDAD_MEDIA], ['20.13%', EFECTIVIDAD_BAJA]    
    ],
    [
        ['20.58%', EFECTIVIDAD_MEDIA], ['23.24%', EFECTIVIDAD_BAJA]    
    ],
    [
        ['10.46%', EFECTIVIDAD_MEDIA], ['40,48%', EFECTIVIDAD_BAJA]    
    ],
]

def srt(grp):
    sort_shape = grp.shapes[0]
    print(sort_shape)
    if sort_shape.has_text_frame:
        return sort_shape.text

prs = Presentation('test.pptx')
for slide in prs.slides:
    group_shapes = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP
    ]
    group_shapes.sort(key=srt)
    for group_shape in group_shapes:
        i = int(group_shape.shapes[0].text) - 1
        group_shape.shapes[0].text = ''
        group_shape.shapes[0].text = dates[i]               
        j = 0
        k = 0
        for oval in [shape for shape in group_shape.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type == MSO_SHAPE.OVAL]:
            oval.text = ''
            p = oval.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(leads_citas[j][i])
            font = run.font
            print(oval.text)
            j = j+1

        for textbox in [shape for shape in group_shape.shapes if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape != group_shape.shapes[0]]:
            textbox.text = ''
            print(textbox.text)
            p = textbox.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = efectividades[i][k][0]
            font = run.font
            font.color.rgb = efectividades[i][k][1]
            print(textbox.text)
            k = k+1

        print('---------')
prs.save('reporte_' + date.today().strftime("%d_%m_%Y") + '.pptx')

