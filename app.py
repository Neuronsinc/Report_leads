from cgitb import text
import datetime
from io import BytesIO
import io
import math
import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from datetime import date
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import requests

# Colores para indicadores de efectividad
EFECTIVIDAD_BAJA = RGBColor(0xFF, 0x00, 0x00)
EFECTIVIDAD_MEDIA = RGBColor(0xFF, 0xA5, 0x00)
EFECTIVIDAD_ALTA = RGBColor(0x00, 0xFF, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_PURPLE = RGBColor(0x56, 0x41, 0x65)
GRAY = RGBColor(0x6D, 0x6D, 0x6D)

# Funcion para sort de grupos de formas
def srt(grp):
    sort_shape = grp.shapes[0]
    if sort_shape.has_text_frame:
        if sort_shape.text.isnumeric():
            return int(sort_shape.text)
    return 100

# Funcion para sort de formas dentro de un grupo
def sort_shape_in_group(shape):

    if (
        shape.has_text_frame
        and shape.text.replace("%", "").isnumeric()
        and int(shape.text.replace("%", "")) > 0
    ):
        return int(shape.text.replace("%", ""))
    elif shape.has_text_frame and (
        shape.text == "total" or shape.text == "pasado" or shape.text == "costos"
    ):
        return 0
    else:
        return 100

# Funcion para generar archivo pptx
def generate_pptx(prs):
    print("-------")

    # Verificar: 
    # que se tienen pares de efectividades por cada semana
    if (
        len(efectividades) == number_of_weeks
        and all([len(x) == 2 for x in efectividades[0]])
        and all(x > 0 for x in st.session_state.projects[proyecto]["matriz_leads_citas"][i])):

        for slide in prs.slides:
            group_shapes = [
                shape
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP
            ]

            group_shapes.sort(key=srt)
            for group_shape in group_shapes:

                shapes_in_group = [shape for shape in group_shape.shapes]
                shapes_in_group.sort(key=sort_shape_in_group)

                if (
                    shapes_in_group[0].has_text_frame
                    and shapes_in_group[0].text.isnumeric()
                    and int(shapes_in_group[0].text) <= 5
                ):
                    n = int(shapes_in_group[0].text) - 1
                    if n < number_of_weeks:
                        shapes_in_group[0].text = ""
                        frame = shapes_in_group[0].text_frame.paragraphs[0]
                        frame.alignment = PP_ALIGN.LEFT
                        run = frame.add_run()
                        run.text = dates[n]
                        font = run.font
                        font.name = "Helvetica Neue"
                        font.size = Pt(14)
                        font.bold = True
                        font.color.rgb = DARK_PURPLE
                        j = 0
                        k = 0
                        for oval in [
                            shape
                            for shape in shapes_in_group
                            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                            and shape.auto_shape_type == MSO_SHAPE.OVAL
                        ]:
                            if j < 3:
                                oval.text = ""
                                frame1 = oval.text_frame.paragraphs[0]
                                frame1.alignment = PP_ALIGN.CENTER
                                run1 = frame1.add_run()
                                run1.text = str(st.session_state.projects[proyecto]["matriz_leads_citas"][n][j])
                                font = run1.font
                                font.name = "Helvetica Neue"
                                font.size = Pt(24 if len(run1.text) == 2 else 18)
                                font.bold = True
                                font.color.rgb = WHITE
                                j = j + 1
                        for textbox in [
                            shape
                            for shape in shapes_in_group
                            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                            and shape != shapes_in_group[0]
                        ]:
                            if k < 2:
                                textbox.text = ""
                                frame2 = textbox.text_frame.paragraphs[0]
                                frame2.alignment = PP_ALIGN.CENTER
                                run2 = frame2.add_run()
                                run2.text = efectividades[n][k][0]
                                font = run2.font
                                font.color.rgb = efectividades[n][k][1]
                                k = k + 1
                elif (
                    shapes_in_group[0].has_text_frame
                    and shapes_in_group[0].text == "total"
                ):
                    shapes_in_group[0].text = ""
                    frame = shapes_in_group[0].text_frame.paragraphs[0]
                    frame.alignment = PP_ALIGN.LEFT
                    run = frame.add_run()
                    run.text = curr_month + " 01" + "-" + today_date.strftime("%d")
                    font = run.font
                    font.name = "Helvetica Neue"
                    font.size = Pt(14)
                    font.bold = True
                    font.color.rgb = DARK_PURPLE
                    m = 0
                    u = 0
                    sub = ["Leads", "Citas Agendadas", "Visitas"]
                    for rectangle in [
                        shape
                        for shape in shapes_in_group
                        if (
                            shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                            and shape != shapes_in_group[0]
                        )
                        or (
                            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                            and shape.auto_shape_type == MSO_SHAPE.RECTANGLE
                        )
                    ]:

                        if "%" in rectangle.text:
                            rectangle.text = ""
                            frame4 = rectangle.text_frame.paragraphs[0]
                            frame4.alignment = PP_ALIGN.CENTER
                            run4 = frame4.add_run()
                            run4.text = total_efectividades[u][0]
                            font = run4.font
                            font.color.rgb = total_efectividades[u][1]
                            u = u + 1
                        if m < 3:
                            rectangle.text = ""
                            frame3 = rectangle.text_frame.paragraphs[0]
                            frame3.alignment = PP_ALIGN.CENTER
                            run3 = frame3.add_run()
                            run3.text = str(totales[m]) + "\n"
                            font = run3.font
                            font.name = "Helvetica Neue"
                            font.size = Pt(20)
                            font.bold = True
                            if m == 1:
                                font.color.rgb = GRAY
                                rectangle.fill.background()
                            else:
                                font.color.rgb = WHITE
                            run32 = frame3.add_run()
                            run32.text = sub[m]
                            font = run32.font
                            font.name = "Helvetica Neue"
                            font.size = Pt(10)
                            font.bold = True
                            font.color.rgb = WHITE if m != 1 else GRAY
                            m = m + 1
                elif (
                    shapes_in_group[0].has_text_frame
                    and shapes_in_group[0].text == "pasado"
                ):
                    shapes_in_group[0].text = ""
                    frame = shapes_in_group[0].text_frame.paragraphs[0]
                    frame.alignment = PP_ALIGN.LEFT
                    run = frame.add_run()
                    run.text = past_month + " 01" + "-" + today_date.strftime("%d")
                    font = run.font
                    font.name = "Helvetica Neue"
                    font.size = Pt(14)
                    font.bold = True
                    font.color.rgb = DARK_PURPLE
                    m = 0
                    u = 0
                    sub = ["Leads", "Citas Agendadas", "Visitas"]
                    for rectangle in [
                        shape
                        for shape in shapes_in_group
                        if (
                            shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                            and shape != shapes_in_group[0]
                        )
                        or (
                            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                            and shape.auto_shape_type == MSO_SHAPE.RECTANGLE
                        )
                    ]:

                        if "%" in rectangle.text:
                            rectangle.text = ""
                            frame4 = rectangle.text_frame.paragraphs[0]
                            frame4.alignment = PP_ALIGN.CENTER
                            run4 = frame4.add_run()
                            run4.text = total_efectividades_pasado[u][0]
                            font = run4.font
                            font.color.rgb = total_efectividades_pasado[u][1]
                            u = u + 1
                        if m < 3:
                            rectangle.text = ""
                            frame3 = rectangle.text_frame.paragraphs[0]
                            frame3.alignment = PP_ALIGN.CENTER
                            run3 = frame3.add_run()
                            run3.text = str(totales_pasado[m]) + "\n"
                            font = run3.font
                            font.name = "Helvetica Neue"
                            font.size = Pt(20)
                            font.bold = True
                            if m == 1:
                                font.color.rgb = GRAY
                                rectangle.fill.background()
                            else:
                                font.color.rgb = WHITE
                            run32 = frame3.add_run()
                            run32.text = sub[m]
                            font = run32.font
                            font.name = "Helvetica Neue"
                            font.size = Pt(10)
                            font.bold = True
                            font.color.rgb = WHITE if m != 1 else GRAY
                            m = m + 1
                elif (
                    shapes_in_group[0].has_text_frame
                    and shapes_in_group[0].text == "costos"
                ):
                    shapes_in_group[0].text = ""
                    m = 0
                    for rectangle in [
                        shape
                        for shape in shapes_in_group
                        if (
                            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                            and shape.auto_shape_type == MSO_SHAPE.RECTANGLE
                        )
                    ]:
                        if m < 2:
                            rectangle.text = ""
                            frame3 = rectangle.text_frame.paragraphs[0]
                            frame3.alignment = PP_ALIGN.CENTER
                            run3 = frame3.add_run()
                            run3.text = str(precios[m])
                            font = run3.font
                            font.name = "Helvetica Neue"
                            font.size = Pt(20)
                            font.bold = True
                            font.color.rgb = WHITE
                            run32 = frame3.add_run()
                            run32.text = " USD"
                            font = run32.font
                            font.name = "Helvetica Neue"
                            font.size = Pt(10)
                            font.bold = True
                            font.color.rgb = WHITE
                            m = m + 1


        binary_output = BytesIO()
        prs.save(binary_output)
        return binary_output.getvalue()


def set_colors(efectividad_agendada, efectividad_realizada, lista_efectividad):
    if efectividad_agendada < 7:
        lista_efectividad.append(
            [str(format(efectividad_agendada, ".2f")) + "%", EFECTIVIDAD_BAJA]
        )
    elif efectividad_agendada >= 7 and efectividad_agendada < 10:
        lista_efectividad.append(
            [str(format(efectividad_agendada, ".2f")) + "%", EFECTIVIDAD_MEDIA]
        )
    elif efectividad_agendada >= 10:
        lista_efectividad.append(
            [str(format(efectividad_agendada, ".2f")) + "%", EFECTIVIDAD_ALTA]
        )

    if efectividad_realizada < 5:
        lista_efectividad.append(
            [str(format(efectividad_realizada, ".2f")) + "%", EFECTIVIDAD_BAJA]
        )
    elif efectividad_realizada >= 5 and efectividad_realizada < 6:
        lista_efectividad.append(
            [str(format(efectividad_realizada, ".2f")) + "%", EFECTIVIDAD_MEDIA]
        )
    elif efectividad_realizada >= 6:
        lista_efectividad.append(
            [str(format(efectividad_realizada, ".2f")) + "%", EFECTIVIDAD_ALTA]
        )


st.set_page_config(layout="wide")
st.title("Creación de Reportes Sobre Leads y Efectividad de Citas")

if "projects" not in st.session_state:
    st.session_state.projects = [
    {
        "start_date"        : "",
        "matriz_leads_citas": [],
        "actual_leads"      : 0,
        "actual_agendada"   : 0,
        "actual_realizada"  : 0,
        "total_leads"       : 0,
        "total_agendada"    : 0,
        "total_realizada"   : 0,
        "pasado_leads"      : 0,
        "pasado_agendada"   : 0,
        "pasado_realizada"  : 0,
        "precio_lead_actual": 0,
        "precio_lead_pasado": 0
    }
]

number_of_weeks = 0
for proyecto in range(len(st.session_state.projects)):
    index = str(proyecto)
    name = st.text_input(label="Nombre del proyecto", key='name ' + index)
    with st.expander(name if name != "" else "proyecto " + index):
        st.session_state.projects[proyecto]["start_date"] = st.date_input(
            "Ingrese fecha inicial:", key="start_date" + index 
        )
        dates = []
        mydate = datetime.datetime.now()
        curr_month = st.session_state.projects[proyecto]["start_date"].strftime("%B").upper()[0:3]
        past_month = (mydate - datetime.timedelta(days=31)).strftime("%B").upper()[0:3]
        primer_dia_mes = datetime.date.today().replace(day=1)

        today_date = datetime.date.today()
        days = abs(today_date - st.session_state.projects[proyecto]["start_date"]).days

        number_of_weeks = math.ceil(days / 7)

        st.session_state.projects[proyecto]["matriz_leads_citas"] = [[] for x in range(number_of_weeks)]

        col1, col2, col3 = st.columns(3)
        if(number_of_weeks > 0):
            col1.write("Numero de Leads")
            col2.write("Citas Agendadas")
            col3.write("Citas Realizadas")

        curr_date = st.session_state.projects[proyecto]["start_date"]
        for i in range(number_of_weeks):
            col1.write(str(curr_date))
            col3.markdown(
                '<p style="color:rgba(0, 0, 0, 0.5)">.</p>', unsafe_allow_html=True
            )
            if curr_date + datetime.timedelta(days=6) > today_date:
                dates.append(
                    curr_month
                    + " "
                    + curr_date.strftime("%d")
                    + "-"
                    + today_date.strftime("%d")
                )
                col2.write(today_date.strftime("%d/%m/%Y"))
            else:
                dates.append(
                    curr_month
                    + " "
                    + curr_date.strftime("%d")
                    + "-"
                    + (curr_date + datetime.timedelta(days=6)).strftime("%d")
                )
                col2.write(str(curr_date + datetime.timedelta(days=6)))
                curr_date = curr_date + datetime.timedelta(days=7)

            st.session_state.projects[proyecto]["matriz_leads_citas"][i].append(
                col1.number_input(
                    value=0,
                    key="Lead" + index + "-" + str(i),
                    min_value=0,
                    label="leads",
                    label_visibility="hidden",
                )
            )
            st.session_state.projects[proyecto]["matriz_leads_citas"][i].append(
                col2.number_input(
                    value=0,
                    key="citaA" + index + "-" + str(i),
                    min_value=0,
                    label="agendadas",
                    label_visibility="hidden",
                )
            )
            st.session_state.projects[proyecto]["matriz_leads_citas"][i].append(
                col3.number_input(
                    value=0,
                    key="citaR" + index + "-" + str(i),
                    min_value=0,
                    label="realizadas",
                    label_visibility="hidden",
                )
            )

        column1, column2 = st.columns(2)
        with column1:
            st.header("Datos del mes actual")
            

            st.session_state.projects[proyecto]["actual_leads"] = st.number_input(value=0, label="Numero de leads mes actual", key="actual_leads" + index)
            st.session_state.projects[proyecto]["actual_agendada"] = st.number_input(
                value=0, label="Numero de citas agendadas mes actual", key="actual_agendada" + index
            )
            st.session_state.projects[proyecto]["actual_realizada"] = st.number_input(
                value=0, label="Numero de citas realizadas mes actual", key="actual_realizada" + index
            )

            st.session_state.projects[proyecto]["total_leads"] = (
                st.session_state.projects[proyecto]["actual_leads"]
                if st.session_state.projects[proyecto]["actual_leads"] > 0
                else sum(l[0] for l in st.session_state.projects[proyecto]["matriz_leads_citas"])
            )
            st.session_state.projects[proyecto]["total_agendada"] = (
                st.session_state.projects[proyecto]["actual_agendada"]
                if st.session_state.projects[proyecto]["actual_agendada"] > 0
                else sum(l[1] for l in st.session_state.projects[proyecto]["matriz_leads_citas"])
            )
            st.session_state.projects[proyecto]["total_realizada"] = (
                st.session_state.projects[proyecto]["actual_realizada"]
                if st.session_state.projects[proyecto]["actual_realizada"] > 0
                else sum(l[2] for l in st.session_state.projects[proyecto]["matriz_leads_citas"])
            )
            totales = [st.session_state.projects[proyecto]["total_leads"], st.session_state.projects[proyecto]["total_agendada"], st.session_state.projects[proyecto]["total_realizada"]]
        with column2:
            st.header("Datos del mes pasado")
            st.session_state.projects[proyecto]["pasado_leads"] = st.number_input(value=0, label="Numero de leads mes pasado", key="pasado_leads" + index)
            st.session_state.projects[proyecto]["pasado_agendada"] = st.number_input(
                value=0, label="Numero de citas agendadas mes pasado", key="pasado_agendada" + index
            )
            st.session_state.projects[proyecto]["pasado_realizada"] = st.number_input(
                value=0, label="Numero de citas realizadas mes pasado", key="pasado_realizada" + index
            )
            totales_pasado = [st.session_state.projects[proyecto]["pasado_leads"], st.session_state.projects[proyecto]["pasado_agendada"], st.session_state.projects[proyecto]["pasado_realizada"]]

            total_efectividades = []
            total_efectividades_pasado = []
 
            if st.session_state.projects[proyecto]["total_agendada"] > 0 and st.session_state.projects[proyecto]["total_leads"] > 0 and st.session_state.projects[proyecto]["total_realizada"] > 0:
                efectividad_agendada_total = st.session_state.projects[proyecto]["total_agendada"] / st.session_state.projects[proyecto]["total_leads"] * 100
                efectividad_realizada_total = st.session_state.projects[proyecto]["total_realizada"] / st.session_state.projects[proyecto]["total_agendada"] * 100
                set_colors(
                    efectividad_agendada_total,
                    efectividad_realizada_total,
                    total_efectividades,
                )
            if st.session_state.projects[proyecto]["pasado_agendada"] > 0 and st.session_state.projects[proyecto]["pasado_leads"] > 0 and st.session_state.projects[proyecto]["pasado_realizada"] > 0:
                efectividad_agendada_pasado = st.session_state.projects[proyecto]["pasado_agendada"] / st.session_state.projects[proyecto]["pasado_leads"] * 100
                efectividad_realizada_pasado = st.session_state.projects[proyecto]["pasado_realizada"] / st.session_state.projects[proyecto]["pasado_agendada"] * 100
                set_colors(
                    efectividad_agendada_pasado,
                    efectividad_realizada_pasado,
                    total_efectividades_pasado,
                )

            show = False
            for efec in total_efectividades_pasado:
                if all(x != "" for x in efec):
                    show = True

        """
        ## Otros
        """
        st.session_state.projects[proyecto]["precio_lead_actual"] = st.number_input(
            value=0.0, label="Precio por lead mes actual", min_value=0.0, key="precio_actual" + index
        )
        st.session_state.projects[proyecto]["precio_lead_pasado"] = st.number_input(
            value=0.0, label="Precio por lead mes pasado", min_value=0.0, key="precio_pasado" + index
        )

        precios = [st.session_state.projects[proyecto]["precio_lead_actual"], st.session_state.projects[proyecto]["precio_lead_pasado"]]

    if proyecto+1 == len(st.session_state.projects):
        if st.button("Agregar Proyecto"):
            st.session_state.projects.append({
        "start_date"        : "",
        "matriz_leads_citas": [],
        "actual_leads"      : 0,
        "actual_agendada"   : 0,
        "actual_realizada"  : 0,
        "total_leads"       : 0,
        "total_agendada"    : 0,
        "total_realizada"   : 0,
        "pasado_leads"      : 0,
        "pasado_agendada"   : 0,
        "pasado_realizada"  : 0,
        "precio_lead_actual": 0,
        "precio_lead_pasado": 0
    })

            st.experimental_rerun()


    if number_of_weeks > 0 and (all(x > 0 for x in st.session_state.projects[proyecto]["matriz_leads_citas"][i])):
        efectividades = [[] for x in range(number_of_weeks)]

        efectividad_agendada = []
        efectividad_realizada = []
        color_ef_agendada = []
        color_ef_realizada = []

        for i in range(number_of_weeks):
            if all(x > 0 for x in st.session_state.projects[proyecto]["matriz_leads_citas"][i]):
                efectividad_agendada = (
                    round(st.session_state.projects[proyecto]["matriz_leads_citas"][i][1] / st.session_state.projects[proyecto]["matriz_leads_citas"][i][0], 4) * 100
                )  # (citas_agendadas / leads) * 100
                efectividad_realizada = (
                    round(st.session_state.projects[proyecto]["matriz_leads_citas"][i][2] / st.session_state.projects[proyecto]["matriz_leads_citas"][i][1], 4) * 100
                )  # (citas_realizadas / citas_agendadas) * 100
                set_colors(efectividad_agendada, efectividad_realizada, efectividades[i])


        if show:
            r = requests.get(
                "https://github.com/Neuronsinc/Report_leads/blob/main/template.pptx?raw=true"
            )
            prs = Presentation(io.BytesIO(r.content))


            
            st.download_button(
                label="Descargar pptx",
                data=generate_pptx(prs),
                file_name="efectividad_" + date.today().strftime("%d_%m_%Y") + ".pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        else:
            st.warning("Debe llenar los datos del mes anterior")
